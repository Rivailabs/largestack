"""PPTX + XLSX document loaders (v0.13.0).

Indian fintechs send Excel-based loan books and PowerPoint pitch decks
constantly. v0.12 said "convert to PDF first" — v0.13 reads them
directly.

Both loaders return the standard LARGESTACK format:

    [{"content": str, "metadata": dict}, ...]

Each slide / sheet becomes one document.

Optional dependencies:

- ``python-pptx`` for ``load_pptx``
- ``openpyxl`` for ``load_xlsx``

If neither is installed, ``load_*`` raises ``ImportError`` with
install hint.
"""

from __future__ import annotations
import asyncio
import logging
from pathlib import Path
from typing import Any

log = logging.getLogger("largestack.loaders.office")


# -------------------- PPTX --------------------


def _have_pptx() -> bool:
    try:
        import pptx  # noqa

        return True
    except ImportError:
        return False


async def load_pptx(path: str | Path) -> list[dict[str, Any]]:
    """Load a .pptx file. One document per slide.

    Each slide's text is concatenated across all shapes (title +
    bullets + tables). Speaker notes go into ``metadata['notes']``.
    """
    if not _have_pptx():
        raise ImportError(
            "python-pptx required for load_pptx. Install with: pip install python-pptx"
        )

    p = Path(path)
    if not p.exists():
        raise FileNotFoundError(f"file not found: {p}")

    return await asyncio.to_thread(_load_pptx_sync, p)


def _load_pptx_sync(path: Path) -> list[dict[str, Any]]:
    from pptx import Presentation

    prs = Presentation(str(path))
    docs: list[dict[str, Any]] = []

    for idx, slide in enumerate(prs.slides, start=1):
        text_parts: list[str] = []
        title = ""

        # Walk shapes
        for shape in slide.shapes:
            # Slide title
            if shape.has_text_frame:
                tf_text = "\n".join(p.text for p in shape.text_frame.paragraphs).strip()
                if not tf_text:
                    continue
                # First text frame in placeholder is usually the title
                if (
                    not title
                    and getattr(shape, "is_placeholder", False)
                    and getattr(
                        shape.placeholder_format,
                        "idx",
                        -1,
                    )
                    == 0
                ):
                    title = tf_text
                text_parts.append(tf_text)

            # Tables on slides
            if shape.has_table:
                rows = []
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    rows.append(" | ".join(cells))
                if rows:
                    text_parts.append("\n".join(rows))

        # Speaker notes
        notes = ""
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            notes = "\n".join(p.text for p in notes_tf.paragraphs).strip()

        content = "\n\n".join(t for t in text_parts if t).strip()
        docs.append(
            {
                "content": content,
                "metadata": {
                    "source": str(path),
                    "loader": "pptx",
                    "slide_number": idx,
                    "title": title,
                    "notes": notes,
                },
            }
        )

    return docs


# -------------------- XLSX --------------------


def _have_openpyxl() -> bool:
    try:
        import openpyxl  # noqa

        return True
    except ImportError:
        return False


async def load_xlsx(
    path: str | Path,
    *,
    rows_per_doc: int | None = None,
) -> list[dict[str, Any]]:
    """Load an .xlsx file. One document per sheet by default.

    Args:
        rows_per_doc: if set, splits each sheet into multiple documents
            of N rows each. Useful for very large sheets.
    """
    if not _have_openpyxl():
        raise ImportError("openpyxl required for load_xlsx. Install with: pip install openpyxl")

    p = Path(path)
    if not p.exists():
        raise FileNotFoundError(f"file not found: {p}")

    return await asyncio.to_thread(_load_xlsx_sync, p, rows_per_doc)


def _load_xlsx_sync(
    path: Path,
    rows_per_doc: int | None,
) -> list[dict[str, Any]]:
    import openpyxl

    wb = openpyxl.load_workbook(str(path), data_only=True, read_only=True)
    docs: list[dict[str, Any]] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))

        if not rows:
            docs.append(
                {
                    "content": "",
                    "metadata": {
                        "source": str(path),
                        "loader": "xlsx",
                        "sheet": sheet_name,
                        "row_count": 0,
                    },
                }
            )
            continue

        # First non-empty row as header
        header: list[str] = []
        data_rows: list[tuple] = []
        for row in rows:
            if any(v not in (None, "") for v in row):
                if not header:
                    header = [str(v) if v is not None else "" for v in row]
                else:
                    data_rows.append(row)

        def render_chunk(chunk: list[tuple]) -> str:
            lines: list[str] = []
            if header:
                lines.append(" | ".join(header))
                lines.append("-" * min(120, len(lines[0])))
            for row in chunk:
                cells = ["" if v is None else str(v) for v in row]
                # Pad/truncate to header length
                if header and len(cells) < len(header):
                    cells = cells + [""] * (len(header) - len(cells))
                lines.append(" | ".join(cells))
            return "\n".join(lines)

        if rows_per_doc is None or len(data_rows) <= rows_per_doc:
            docs.append(
                {
                    "content": render_chunk(data_rows),
                    "metadata": {
                        "source": str(path),
                        "loader": "xlsx",
                        "sheet": sheet_name,
                        "row_count": len(data_rows),
                        "header": header,
                    },
                }
            )
        else:
            for i in range(0, len(data_rows), rows_per_doc):
                chunk = data_rows[i : i + rows_per_doc]
                docs.append(
                    {
                        "content": render_chunk(chunk),
                        "metadata": {
                            "source": str(path),
                            "loader": "xlsx",
                            "sheet": sheet_name,
                            "row_count": len(chunk),
                            "row_offset": i,
                            "header": header,
                        },
                    }
                )

    wb.close()
    return docs


__all__ = ["load_pptx", "load_xlsx"]
