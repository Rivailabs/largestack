"""Document loaders for common file formats (v0.7.0).

A document loader takes a file path or file-like object and returns a
list of dicts: ``[{"content": str, "metadata": dict}, ...]``. This
output format matches LARGESTACK's RAG ingestion pipeline directly.

Loaders included:
- ``load_text(path)`` — plain .txt/.text/anything
- ``load_markdown(path)`` — .md (parses YAML frontmatter if present)
- ``load_pdf(path)`` — .pdf via pypdf (one document per page)
- ``load_docx(path)`` — .docx via python-docx
- ``load_html(path_or_url)`` — .html via beautifulsoup4 (cleans tags)
- ``load_csv(path)`` — .csv (one document per row)
- ``load_json(path)`` — .json
- ``load_yaml(path)`` — .yaml / .yml
- ``load_xml(path)`` — .xml

All loaders return the same schema regardless of source format. Use
the dispatcher ``load(path)`` to auto-detect format from extension.

Optional dependencies (loaders gracefully report when missing):
- pypdf for PDF
- python-docx for DOCX
- beautifulsoup4 + httpx for HTML
- pyyaml for YAML
"""
from __future__ import annotations
import asyncio
import csv
import json
import logging
import os
import re
from defusedxml import ElementTree as ET
from defusedxml.common import DefusedXmlException
from typing import Any

log = logging.getLogger("largestack.loaders")


# -------------------- Plain text --------------------

async def load_text(path: str, encoding: str = "utf-8") -> list[dict]:
    """Load a plain-text file as a single document.

    Args:
        path: file path
        encoding: text encoding (default utf-8). Falls back to latin-1
            if utf-8 fails.

    Returns:
        ``[{"content": str, "metadata": {...}}]``
    """
    try:
        async def _read():
            try:
                with open(path, "r", encoding=encoding) as f:
                    return f.read()
            except UnicodeDecodeError:
                with open(path, "r", encoding="latin-1") as f:
                    return f.read()
        text = await asyncio.to_thread(lambda: _read_sync(path, encoding))
    except FileNotFoundError:
        return [{"content": "", "metadata": {"error": f"file not found: {path}"}}]
    except Exception as e:
        return [{"content": "", "metadata": {"error": str(e)}}]
    return [{"content": text, "metadata": {"source": path, "format": "text"}}]


def _read_sync(path: str, encoding: str) -> str:
    try:
        with open(path, "r", encoding=encoding) as f:
            return f.read()
    except UnicodeDecodeError:
        with open(path, "r", encoding="latin-1") as f:
            return f.read()


# -------------------- Markdown --------------------

_FRONTMATTER_RE = re.compile(r"^---\s*\n(.*?)\n---\s*\n", re.DOTALL)


async def load_markdown(path: str) -> list[dict]:
    """Load a markdown file. Strips YAML frontmatter into metadata.

    Returns the markdown body as content; frontmatter is parsed (if pyyaml
    is installed) and merged into metadata.
    """
    try:
        text = await asyncio.to_thread(_read_sync, path, "utf-8")
    except FileNotFoundError:
        return [{"content": "", "metadata": {"error": f"file not found: {path}"}}]

    metadata: dict[str, Any] = {"source": path, "format": "markdown"}

    # Extract YAML frontmatter if present
    m = _FRONTMATTER_RE.match(text)
    if m:
        fm_raw = m.group(1)
        body = text[m.end():]
        try:
            import yaml  # type: ignore
            fm = yaml.safe_load(fm_raw)
            if isinstance(fm, dict):
                metadata.update(fm)
        except ImportError:
            metadata["frontmatter_raw"] = fm_raw  # leave as raw if no yaml
        except Exception as e:
            log.debug(f"frontmatter parse failed: {e}")
        text = body

    return [{"content": text, "metadata": metadata}]


# -------------------- PDF --------------------

async def load_pdf(path: str) -> list[dict]:
    """Load a PDF, one document per page.

    Requires ``pip install pypdf``. Without it, returns an error doc.
    """
    try:
        from pypdf import PdfReader
    except ImportError:
        return [{
            "content": "",
            "metadata": {"error": "PDF loader needs: pip install pypdf"},
        }]

    def _extract():
        reader = PdfReader(path)
        return [(i, p.extract_text() or "") for i, p in enumerate(reader.pages)]

    try:
        pages = await asyncio.to_thread(_extract)
    except FileNotFoundError:
        return [{"content": "", "metadata": {"error": f"file not found: {path}"}}]
    except Exception as e:
        return [{"content": "", "metadata": {"error": f"PDF read failed: {e}"}}]

    return [
        {
            "content": text,
            "metadata": {"source": path, "format": "pdf", "page": i, "total_pages": len(pages)},
        }
        for i, text in pages
    ]


# -------------------- DOCX --------------------

async def load_docx(path: str) -> list[dict]:
    """Load a .docx file as a single document (paragraphs joined by newline).

    Requires ``pip install python-docx``.
    """
    try:
        import docx  # type: ignore  (python-docx)
    except ImportError:
        return [{
            "content": "",
            "metadata": {"error": "DOCX loader needs: pip install python-docx"},
        }]

    def _read():
        d = docx.Document(path)
        paras = [p.text for p in d.paragraphs if p.text]
        return "\n".join(paras)

    try:
        text = await asyncio.to_thread(_read)
    except FileNotFoundError:
        return [{"content": "", "metadata": {"error": f"file not found: {path}"}}]
    except Exception as e:
        return [{"content": "", "metadata": {"error": f"DOCX read failed: {e}"}}]
    return [{"content": text, "metadata": {"source": path, "format": "docx"}}]


# -------------------- HTML --------------------

async def load_html(path_or_url: str, *, fetch_remote: bool = True) -> list[dict]:
    """Load an HTML file or URL. Strips tags, returns clean text.

    If ``path_or_url`` starts with http:// or https:// AND
    ``fetch_remote=True``, fetches via httpx. Otherwise reads from disk.

    Requires ``pip install beautifulsoup4`` (and httpx for URLs — already
    a LARGESTACK dep).
    """
    try:
        from bs4 import BeautifulSoup  # type: ignore
    except ImportError:
        return [{
            "content": "",
            "metadata": {"error": "HTML loader needs: pip install beautifulsoup4"},
        }]

    is_url = path_or_url.startswith(("http://", "https://"))
    if is_url and fetch_remote:
        try:
            import httpx
            async with httpx.AsyncClient(timeout=15) as c:
                r = await c.get(path_or_url, follow_redirects=True)
                if r.status_code >= 400:
                    return [{
                        "content": "",
                        "metadata": {"error": f"HTTP {r.status_code} fetching {path_or_url}"},
                    }]
                raw = r.text
        except Exception as e:
            return [{"content": "", "metadata": {"error": f"HTTP fetch failed: {e}"}}]
    else:
        try:
            raw = await asyncio.to_thread(_read_sync, path_or_url, "utf-8")
        except FileNotFoundError:
            return [{"content": "", "metadata": {"error": f"file not found: {path_or_url}"}}]

    def _parse():
        soup = BeautifulSoup(raw, "html.parser")
        # Remove script/style/nav/footer
        for tag in soup(["script", "style", "nav", "footer", "aside"]):
            tag.decompose()
        title_el = soup.find("title")
        title = title_el.get_text(strip=True) if title_el else ""
        text = soup.get_text(separator="\n", strip=True)
        # Collapse multiple blank lines
        text = re.sub(r"\n\s*\n+", "\n\n", text)
        return title, text

    try:
        title, text = await asyncio.to_thread(_parse)
    except Exception as e:
        return [{"content": "", "metadata": {"error": f"HTML parse failed: {e}"}}]

    return [{
        "content": text,
        "metadata": {
            "source": path_or_url,
            "format": "html",
            "title": title,
            "is_url": is_url,
        },
    }]


# -------------------- CSV --------------------

async def load_csv(path: str, *, has_header: bool = True) -> list[dict]:
    """Load a CSV file, one document per row.

    Each row's content is its dict-like representation; metadata
    includes the row index and source path.
    """
    def _read():
        rows = []
        with open(path, "r", encoding="utf-8", newline="") as f:
            if has_header:
                reader = csv.DictReader(f)
                for i, row in enumerate(reader):
                    content = "\n".join(f"{k}: {v}" for k, v in row.items())
                    rows.append((i, content, dict(row)))
            else:
                reader = csv.reader(f)
                for i, row in enumerate(reader):
                    content = " | ".join(str(c) for c in row)
                    rows.append((i, content, {"raw_row": row}))
        return rows

    try:
        rows = await asyncio.to_thread(_read)
    except FileNotFoundError:
        return [{"content": "", "metadata": {"error": f"file not found: {path}"}}]
    except Exception as e:
        return [{"content": "", "metadata": {"error": f"CSV read failed: {e}"}}]

    return [
        {
            "content": content,
            "metadata": {
                "source": path,
                "format": "csv",
                "row": i,
                "fields": fields,
            },
        }
        for i, content, fields in rows
    ]


# -------------------- JSON --------------------

async def load_json(path: str) -> list[dict]:
    """Load a JSON file. Top-level object → 1 doc. Top-level array → N docs.

    JSONL format (one object per line) is also supported via ``load_jsonl()``.
    """
    try:
        text = await asyncio.to_thread(_read_sync, path, "utf-8")
    except FileNotFoundError:
        return [{"content": "", "metadata": {"error": f"file not found: {path}"}}]

    try:
        data = json.loads(text)
    except json.JSONDecodeError as e:
        return [{"content": "", "metadata": {"error": f"JSON parse error: {e}"}}]

    if isinstance(data, list):
        return [
            {
                "content": json.dumps(item, indent=2),
                "metadata": {"source": path, "format": "json", "index": i},
            }
            for i, item in enumerate(data)
        ]
    return [{
        "content": json.dumps(data, indent=2),
        "metadata": {"source": path, "format": "json"},
    }]


async def load_jsonl(path: str) -> list[dict]:
    """Load a JSONL file (one JSON object per line)."""
    try:
        text = await asyncio.to_thread(_read_sync, path, "utf-8")
    except FileNotFoundError:
        return [{"content": "", "metadata": {"error": f"file not found: {path}"}}]

    docs = []
    for i, line in enumerate(text.splitlines()):
        line = line.strip()
        if not line:
            continue
        try:
            obj = json.loads(line)
            docs.append({
                "content": json.dumps(obj, indent=2),
                "metadata": {"source": path, "format": "jsonl", "line": i},
            })
        except json.JSONDecodeError:
            continue
    return docs


# -------------------- YAML --------------------

async def load_yaml(path: str) -> list[dict]:
    """Load a YAML file as a single document.

    Requires ``pip install pyyaml``.
    """
    try:
        import yaml  # type: ignore
    except ImportError:
        return [{
            "content": "",
            "metadata": {"error": "YAML loader needs: pip install pyyaml"},
        }]

    try:
        text = await asyncio.to_thread(_read_sync, path, "utf-8")
    except FileNotFoundError:
        return [{"content": "", "metadata": {"error": f"file not found: {path}"}}]

    try:
        data = yaml.safe_load(text)
    except Exception as e:
        return [{"content": "", "metadata": {"error": f"YAML parse error: {e}"}}]

    return [{
        "content": yaml.safe_dump(data, default_flow_style=False),
        "metadata": {"source": path, "format": "yaml"},
    }]


# -------------------- XML --------------------

async def load_xml(path: str) -> list[dict]:
    """Load an XML file. Returns the text content with tags preserved."""
    try:
        text = await asyncio.to_thread(_read_sync, path, "utf-8")
    except FileNotFoundError:
        return [{"content": "", "metadata": {"error": f"file not found: {path}"}}]

    # Validate XML; extract text-only version
    try:
        root = ET.fromstring(text)
        text_only = "\n".join(
            (e.text or "").strip() for e in root.iter() if (e.text or "").strip()
        )
        return [{
            "content": text,  # full XML
            "metadata": {
                "source": path,
                "format": "xml",
                "root_tag": root.tag,
                "text_only": text_only,
            },
        }]
    except (ET.ParseError, DefusedXmlException) as e:
        return [{"content": "", "metadata": {"error": f"XML parse error: {e}"}}]


# -------------------- Dispatcher --------------------

async def load(path: str) -> list[dict]:
    """Auto-detect format from extension and load.

    Args:
        path: file path or URL.

    Returns:
        Same dict-list shape as the underlying loader.

    Falls back to ``load_text`` for unknown extensions.
    """
    if path.startswith(("http://", "https://")):
        return await load_html(path)

    ext = os.path.splitext(path)[1].lower().lstrip(".")
    dispatch = {
        "txt": load_text, "text": load_text, "log": load_text,
        "md": load_markdown, "markdown": load_markdown,
        "pdf": load_pdf,
        "docx": load_docx,
        "html": load_html, "htm": load_html,
        "csv": load_csv,
        "json": load_json,
        "jsonl": load_jsonl, "ndjson": load_jsonl,
        "yaml": load_yaml, "yml": load_yaml,
        "xml": load_xml,
        # v0.8.0 additions
        "pptx": load_pptx,
        "epub": load_epub,
        "xlsx": load_excel, "xls": load_excel,
    }
    fn = dispatch.get(ext, load_text)
    return await fn(path)


# -------------------- v0.8.0 New Loaders --------------------

async def load_pptx(path: str) -> list[dict]:
    """Load .pptx file. One document per slide.

    Requires ``pip install python-pptx``.
    """
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        return [{
            "content": "",
            "metadata": {"error": "PPTX loader needs: pip install python-pptx"},
        }]

    def _read():
        prs = Presentation(path)
        slides = []
        for i, slide in enumerate(prs.slides):
            text_parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = "".join(run.text for run in para.runs)
                        if line.strip():
                            text_parts.append(line)
            slides.append((i, "\n".join(text_parts)))
        return slides

    try:
        slides = await asyncio.to_thread(_read)
    except FileNotFoundError:
        return [{"content": "", "metadata": {"error": f"file not found: {path}"}}]
    except Exception as e:
        return [{"content": "", "metadata": {"error": f"PPTX read failed: {e}"}}]

    return [
        {
            "content": text,
            "metadata": {
                "source": path, "format": "pptx",
                "slide": i, "total_slides": len(slides),
            },
        }
        for i, text in slides
    ]


async def load_epub(path: str) -> list[dict]:
    """Load .epub file. One document per chapter/section.

    Requires ``pip install ebooklib``.
    """
    try:
        from ebooklib import epub, ITEM_DOCUMENT  # type: ignore
        from bs4 import BeautifulSoup  # type: ignore
    except ImportError:
        return [{
            "content": "",
            "metadata": {"error": "EPUB loader needs: pip install ebooklib beautifulsoup4"},
        }]

    def _read():
        book = epub.read_epub(path)
        title = book.get_metadata("DC", "title")
        title = title[0][0] if title else ""
        chapters = []
        for i, item in enumerate(book.get_items()):
            if item.get_type() == ITEM_DOCUMENT:
                soup = BeautifulSoup(item.get_content(), "html.parser")
                for tag in soup(["script", "style"]):
                    tag.decompose()
                text = soup.get_text(separator="\n", strip=True)
                if text.strip():
                    chapters.append((i, text, item.get_name() or ""))
        return title, chapters

    try:
        title, chapters = await asyncio.to_thread(_read)
    except FileNotFoundError:
        return [{"content": "", "metadata": {"error": f"file not found: {path}"}}]
    except Exception as e:
        return [{"content": "", "metadata": {"error": f"EPUB read failed: {e}"}}]

    return [
        {
            "content": text,
            "metadata": {
                "source": path, "format": "epub",
                "book_title": title, "chapter": i, "chapter_name": name,
            },
        }
        for i, text, name in chapters
    ]


async def load_excel(path: str, sheet_name: str | None = None) -> list[dict]:
    """Load .xlsx or .xls. One document per sheet (or just one if sheet_name).

    Each sheet's content is rendered as a tab-separated string.
    Requires ``pip install openpyxl``.
    """
    try:
        from openpyxl import load_workbook  # type: ignore
    except ImportError:
        return [{
            "content": "",
            "metadata": {"error": "Excel loader needs: pip install openpyxl"},
        }]

    def _read():
        wb = load_workbook(path, data_only=True, read_only=True)
        sheets_to_read = [sheet_name] if sheet_name else wb.sheetnames
        result = []
        for sname in sheets_to_read:
            if sname not in wb.sheetnames:
                continue
            ws = wb[sname]
            rows = []
            for row in ws.iter_rows(values_only=True):
                rows.append("\t".join(
                    str(c) if c is not None else "" for c in row
                ))
            result.append((sname, "\n".join(rows), ws.max_row, ws.max_column))
        wb.close()
        return result

    try:
        sheets = await asyncio.to_thread(_read)
    except FileNotFoundError:
        return [{"content": "", "metadata": {"error": f"file not found: {path}"}}]
    except Exception as e:
        return [{"content": "", "metadata": {"error": f"Excel read failed: {e}"}}]

    return [
        {
            "content": text,
            "metadata": {
                "source": path, "format": "xlsx",
                "sheet": sname, "rows": rows, "cols": cols,
            },
        }
        for sname, text, rows, cols in sheets
    ]


async def load_s3(bucket: str, key: str, *, region: str | None = None) -> list[dict]:
    """Load a single object from S3 and dispatch based on extension.

    Args:
        bucket: S3 bucket name.
        key: object key (path within bucket).
        region: AWS region (else from env / boto3 default).

    Auth: standard AWS env vars or instance profile.
    Requires ``pip install boto3``.
    """
    try:
        import boto3  # type: ignore
    except ImportError:
        return [{
            "content": "",
            "metadata": {"error": "S3 loader needs: pip install boto3"},
        }]

    def _read():
        client_kw = {"region_name": region} if region else {}
        s3 = boto3.client("s3", **client_kw)
        obj = s3.get_object(Bucket=bucket, Key=key)
        return obj["Body"].read()

    try:
        data = await asyncio.to_thread(_read)
    except Exception as e:
        return [{"content": "", "metadata": {"error": f"S3 fetch failed: {e}"}}]

    # Save to temp file and dispatch by extension
    import tempfile
    ext = os.path.splitext(key)[1] or ".txt"
    with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
        tmp.write(data)
        tmp_path = tmp.name
    try:
        docs = await load(tmp_path)
        for d in docs:
            d["metadata"]["s3_bucket"] = bucket
            d["metadata"]["s3_key"] = key
            d["metadata"]["source"] = f"s3://{bucket}/{key}"
        return docs
    finally:
        try:
            os.unlink(tmp_path)
        except OSError:
            pass


async def load_gcs(bucket: str, blob_name: str) -> list[dict]:
    """Load a single object from Google Cloud Storage.

    Auth: ``GOOGLE_APPLICATION_CREDENTIALS`` env var or default credentials.
    Requires ``pip install google-cloud-storage``.
    """
    try:
        from google.cloud import storage  # type: ignore
    except ImportError:
        return [{
            "content": "",
            "metadata": {"error": "GCS loader needs: pip install google-cloud-storage"},
        }]

    def _read():
        client = storage.Client()
        b = client.bucket(bucket)
        blob = b.blob(blob_name)
        return blob.download_as_bytes()

    try:
        data = await asyncio.to_thread(_read)
    except Exception as e:
        return [{"content": "", "metadata": {"error": f"GCS fetch failed: {e}"}}]

    import tempfile
    ext = os.path.splitext(blob_name)[1] or ".txt"
    with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
        tmp.write(data)
        tmp_path = tmp.name
    try:
        docs = await load(tmp_path)
        for d in docs:
            d["metadata"]["gcs_bucket"] = bucket
            d["metadata"]["gcs_blob"] = blob_name
            d["metadata"]["source"] = f"gs://{bucket}/{blob_name}"
        return docs
    finally:
        try:
            os.unlink(tmp_path)
        except OSError:
            pass


async def load_azure_blob(
    account_url: str, container: str, blob_name: str,
    *, credential: str | None = None,
) -> list[dict]:
    """Load a single blob from Azure Blob Storage.

    Args:
        account_url: e.g. ``"https://myacct.blob.core.windows.net"``.
        container: container name.
        blob_name: blob path within the container.
        credential: SAS token / account key / None (DefaultAzureCredential).

    Requires ``pip install azure-storage-blob``.
    """
    try:
        from azure.storage.blob.aio import BlobServiceClient  # type: ignore
    except ImportError:
        return [{
            "content": "",
            "metadata": {"error": "Azure loader needs: pip install azure-storage-blob"},
        }]

    try:
        async with BlobServiceClient(
            account_url=account_url, credential=credential
        ) as svc:
            blob = svc.get_blob_client(container=container, blob=blob_name)
            stream = await blob.download_blob()
            data = await stream.readall()
    except Exception as e:
        return [{"content": "", "metadata": {"error": f"Azure fetch failed: {e}"}}]

    import tempfile
    ext = os.path.splitext(blob_name)[1] or ".txt"
    with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
        tmp.write(data)
        tmp_path = tmp.name
    try:
        docs = await load(tmp_path)
        for d in docs:
            d["metadata"]["azure_container"] = container
            d["metadata"]["azure_blob"] = blob_name
            d["metadata"]["source"] = f"{account_url}/{container}/{blob_name}"
        return docs
    finally:
        try:
            os.unlink(tmp_path)
        except OSError:
            pass


async def load_youtube_transcript(
    video_id_or_url: str, *, languages: list[str] | None = None,
) -> list[dict]:
    """Load transcript of a YouTube video.

    Args:
        video_id_or_url: 11-char video ID or full YouTube URL.
        languages: preferred language codes (default ``["en"]``).

    Requires ``pip install youtube-transcript-api``.
    """
    try:
        from youtube_transcript_api import YouTubeTranscriptApi  # type: ignore
    except ImportError:
        return [{
            "content": "",
            "metadata": {"error": "YouTube loader needs: pip install youtube-transcript-api"},
        }]

    # Extract video ID from URL if needed
    vid = video_id_or_url
    if "youtube.com" in vid or "youtu.be" in vid:
        m = re.search(r"(?:v=|youtu\.be/|/embed/)([A-Za-z0-9_-]{11})", vid)
        if m:
            vid = m.group(1)

    langs = languages or ["en"]

    def _fetch():
        return YouTubeTranscriptApi.get_transcript(vid, languages=langs)

    try:
        items = await asyncio.to_thread(_fetch)
    except Exception as e:
        return [{"content": "", "metadata": {"error": f"YouTube fetch failed: {e}"}}]

    text = "\n".join(item.get("text", "") for item in items)
    return [{
        "content": text,
        "metadata": {
            "source": f"https://www.youtube.com/watch?v={vid}",
            "format": "youtube_transcript",
            "video_id": vid,
            "n_segments": len(items),
        },
    }]


async def load_wikipedia(query: str, *, lang: str = "en", sentences: int = 0) -> list[dict]:
    """Search Wikipedia and load the matching article's content.

    Uses Wikipedia's REST API directly via httpx — no external SDK
    required. ``sentences=0`` returns full article extract;
    ``sentences=N`` truncates to first N sentences (Wikipedia API
    feature).
    """
    import urllib.parse
    base = f"https://{lang}.wikipedia.org/w/api.php"
    params: dict = {
        "action": "query",
        "format": "json",
        "prop": "extracts",
        "explaintext": "1",
        "redirects": "1",
        "titles": query,
    }
    if sentences > 0:
        params["exsentences"] = str(sentences)

    try:
        import httpx
        async with httpx.AsyncClient(timeout=20) as c:
            r = await c.get(base, params=params)
            r.raise_for_status()
            data = r.json()
    except Exception as e:
        return [{"content": "", "metadata": {"error": f"Wikipedia fetch failed: {e}"}}]

    pages = (data.get("query") or {}).get("pages") or {}
    if not pages:
        return [{"content": "", "metadata": {"error": "no result"}}]
    # First (and usually only) page
    pid, page = next(iter(pages.items()))
    if pid == "-1":
        return [{"content": "", "metadata": {"error": f"no Wikipedia article for {query!r}"}}]
    extract = page.get("extract", "")
    title = page.get("title", query)
    return [{
        "content": extract,
        "metadata": {
            "source": f"https://{lang}.wikipedia.org/wiki/{urllib.parse.quote(title)}",
            "format": "wikipedia",
            "title": title,
            "language": lang,
        },
    }]


async def load_arxiv(query: str, *, max_results: int = 5) -> list[dict]:
    """Search ArXiv and return matching paper abstracts.

    Uses ArXiv's public Atom API directly via httpx — no SDK required.
    Returns one document per paper with ``content`` = abstract,
    ``metadata`` = title, authors, arxiv_id, pdf_url, published date.
    """
    import urllib.parse
    base = "http://export.arxiv.org/api/query"
    params = {
        "search_query": query,
        "max_results": str(max_results),
        "sortBy": "relevance",
        "sortOrder": "descending",
    }
    try:
        import httpx
        async with httpx.AsyncClient(timeout=20) as c:
            r = await c.get(base, params=params)
            r.raise_for_status()
            text = r.text
    except Exception as e:
        return [{"content": "", "metadata": {"error": f"ArXiv fetch failed: {e}"}}]

    # Parse Atom XML
    try:
        ns = {"atom": "http://www.w3.org/2005/Atom"}
        root = ET.fromstring(text)
    except (ET.ParseError, DefusedXmlException) as e:
        return [{"content": "", "metadata": {"error": f"ArXiv parse error: {e}"}}]

    docs = []
    for entry in root.findall("atom:entry", ns):
        title_el = entry.find("atom:title", ns)
        summary_el = entry.find("atom:summary", ns)
        id_el = entry.find("atom:id", ns)
        published_el = entry.find("atom:published", ns)
        title = (title_el.text or "").strip() if title_el is not None else ""
        summary = (summary_el.text or "").strip() if summary_el is not None else ""
        arxiv_id = (id_el.text or "").strip() if id_el is not None else ""
        published = (published_el.text or "") if published_el is not None else ""
        authors = [
            (a.find("atom:name", ns).text or "")
            for a in entry.findall("atom:author", ns)
            if a.find("atom:name", ns) is not None
        ]
        pdf_url = arxiv_id.replace("/abs/", "/pdf/") if "/abs/" in arxiv_id else ""
        docs.append({
            "content": summary,
            "metadata": {
                "source": arxiv_id,
                "format": "arxiv",
                "title": title,
                "authors": authors,
                "published": published,
                "pdf_url": pdf_url,
            },
        })
    if not docs:
        return [{"content": "", "metadata": {"error": "no ArXiv results"}}]
    return docs


async def load_pubmed(query: str, *, max_results: int = 5) -> list[dict]:
    """Search PubMed and return matching abstracts.

    Uses NCBI E-utilities directly via httpx. Two-step:
    1. esearch to get IDs
    2. efetch to get abstracts
    """
    base = "https://eutils.ncbi.nlm.nih.gov/entrez/eutils"
    try:
        import httpx
        async with httpx.AsyncClient(timeout=30) as c:
            # 1. Get IDs
            r = await c.get(
                f"{base}/esearch.fcgi",
                params={
                    "db": "pubmed",
                    "term": query,
                    "retmax": str(max_results),
                    "retmode": "json",
                },
            )
            r.raise_for_status()
            ids = (r.json().get("esearchresult") or {}).get("idlist") or []
            if not ids:
                return [{"content": "", "metadata": {"error": "no PubMed results"}}]
            # 2. Get abstracts (XML)
            r2 = await c.get(
                f"{base}/efetch.fcgi",
                params={"db": "pubmed", "id": ",".join(ids), "retmode": "xml"},
            )
            r2.raise_for_status()
            xml = r2.text
    except Exception as e:
        return [{"content": "", "metadata": {"error": f"PubMed fetch failed: {e}"}}]

    try:
        root = ET.fromstring(xml)
    except (ET.ParseError, DefusedXmlException) as e:
        return [{"content": "", "metadata": {"error": f"PubMed parse error: {e}"}}]

    docs = []
    for art in root.findall(".//PubmedArticle"):
        pmid_el = art.find(".//PMID")
        title_el = art.find(".//ArticleTitle")
        abstract_parts = [
            (a.text or "")
            for a in art.findall(".//AbstractText")
            if (a.text or "").strip()
        ]
        pmid = pmid_el.text if pmid_el is not None else ""
        title = title_el.text if title_el is not None else ""
        abstract = "\n".join(abstract_parts)
        docs.append({
            "content": abstract,
            "metadata": {
                "source": f"https://pubmed.ncbi.nlm.nih.gov/{pmid}/",
                "format": "pubmed",
                "pmid": pmid,
                "title": title,
            },
        })
    if not docs:
        return [{"content": "", "metadata": {"error": "PubMed: no parseable articles"}}]
    return docs


# -------------------- v0.9.0: 8 high-value loaders --------------------

from largestack._loaders.loaders_v09 import (
    load_notion_database,
    load_confluence,
    load_github_repo,
    load_google_drive,
    load_email_imap,
    load_gmail,
    load_web_scrape,
    load_ocr,
)
