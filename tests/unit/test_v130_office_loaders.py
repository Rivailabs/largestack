"""v0.13.0: Tests for pptx + xlsx loaders."""
from __future__ import annotations

import sys
from pathlib import Path
from unittest.mock import patch

import pytest


# -------------------- Module imports --------------------

def test_office_module_imports_cleanly():
    from largestack._loaders import office
    assert hasattr(office, "load_pptx")
    assert hasattr(office, "load_xlsx")


# -------------------- PPTX --------------------

@pytest.mark.asyncio
async def test_load_pptx_missing_dep_raises():
    from largestack._loaders import office

    with patch.object(office, "_have_pptx", return_value=False):
        with pytest.raises(ImportError, match="python-pptx"):
            await office.load_pptx("/tmp/x.pptx")


@pytest.mark.asyncio
async def test_load_pptx_missing_file(tmp_path):
    pytest.importorskip("pptx")
    from largestack._loaders.office import load_pptx
    with pytest.raises(FileNotFoundError):
        await load_pptx(tmp_path / "nope.pptx")


@pytest.mark.asyncio
async def test_load_pptx_extracts_slide_text(tmp_path):
    """Build a real pptx and round-trip it through the loader."""
    pptx = pytest.importorskip("pptx")
    from pptx import Presentation
    from pptx.util import Inches
    from largestack._loaders.office import load_pptx

    # Build a 2-slide presentation
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]

    slide1 = prs.slides.add_slide(blank_layout)
    txbox = slide1.shapes.add_textbox(
        Inches(1), Inches(1), Inches(5), Inches(1),
    )
    txbox.text_frame.text = "First Slide Title"

    slide2 = prs.slides.add_slide(blank_layout)
    txbox = slide2.shapes.add_textbox(
        Inches(1), Inches(1), Inches(5), Inches(1),
    )
    txbox.text_frame.text = "Second Slide Content"

    out = tmp_path / "deck.pptx"
    prs.save(str(out))

    docs = await load_pptx(out)
    assert len(docs) == 2
    assert "First Slide Title" in docs[0]["content"]
    assert "Second Slide Content" in docs[1]["content"]
    assert docs[0]["metadata"]["loader"] == "pptx"
    assert docs[0]["metadata"]["slide_number"] == 1
    assert docs[1]["metadata"]["slide_number"] == 2


# -------------------- XLSX --------------------

@pytest.mark.asyncio
async def test_load_xlsx_missing_dep_raises():
    from largestack._loaders import office
    with patch.object(office, "_have_openpyxl", return_value=False):
        with pytest.raises(ImportError, match="openpyxl"):
            await office.load_xlsx("/tmp/x.xlsx")


@pytest.mark.asyncio
async def test_load_xlsx_missing_file(tmp_path):
    pytest.importorskip("openpyxl")
    from largestack._loaders.office import load_xlsx
    with pytest.raises(FileNotFoundError):
        await load_xlsx(tmp_path / "nope.xlsx")


@pytest.mark.asyncio
async def test_load_xlsx_one_doc_per_sheet(tmp_path):
    openpyxl = pytest.importorskip("openpyxl")
    from openpyxl import Workbook
    from largestack._loaders.office import load_xlsx

    wb = Workbook()
    ws1 = wb.active
    ws1.title = "Loans"
    ws1.append(["loan_id", "amount", "borrower"])
    ws1.append(["L001", 50000, "Sachith"])
    ws1.append(["L002", 100000, "Sushma"])

    ws2 = wb.create_sheet("Repayments")
    ws2.append(["loan_id", "paid_amount", "date"])
    ws2.append(["L001", 5000, "2026-04-01"])

    out = tmp_path / "book.xlsx"
    wb.save(str(out))

    docs = await load_xlsx(out)
    assert len(docs) == 2

    loans_doc = next(d for d in docs if d["metadata"]["sheet"] == "Loans")
    assert "loan_id" in loans_doc["content"]
    assert "L001" in loans_doc["content"]
    assert "Sachith" in loans_doc["content"]
    assert loans_doc["metadata"]["row_count"] == 2

    rep_doc = next(d for d in docs if d["metadata"]["sheet"] == "Repayments")
    assert rep_doc["metadata"]["row_count"] == 1


@pytest.mark.asyncio
async def test_load_xlsx_chunks_when_rows_per_doc_set(tmp_path):
    pytest.importorskip("openpyxl")
    from openpyxl import Workbook
    from largestack._loaders.office import load_xlsx

    wb = Workbook()
    ws = wb.active
    ws.append(["id", "value"])
    for i in range(10):
        ws.append([i, f"row{i}"])

    out = tmp_path / "big.xlsx"
    wb.save(str(out))

    docs = await load_xlsx(out, rows_per_doc=3)
    # 10 rows / 3 per doc = 4 chunks
    assert len(docs) == 4
    assert docs[0]["metadata"]["row_offset"] == 0
    assert docs[1]["metadata"]["row_offset"] == 3
    assert docs[3]["metadata"]["row_count"] == 1  # last chunk


@pytest.mark.asyncio
async def test_load_xlsx_handles_empty_sheet(tmp_path):
    pytest.importorskip("openpyxl")
    from openpyxl import Workbook
    from largestack._loaders.office import load_xlsx

    wb = Workbook()
    ws = wb.active
    ws.title = "Empty"
    out = tmp_path / "empty.xlsx"
    wb.save(str(out))

    docs = await load_xlsx(out)
    assert len(docs) == 1
    assert docs[0]["content"] == ""
    assert docs[0]["metadata"]["row_count"] == 0


@pytest.mark.asyncio
async def test_load_xlsx_includes_header_in_metadata(tmp_path):
    pytest.importorskip("openpyxl")
    from openpyxl import Workbook
    from largestack._loaders.office import load_xlsx

    wb = Workbook()
    ws = wb.active
    ws.append(["aadhaar_last4", "pan_masked", "amount"])
    ws.append(["9012", "AAA***1C", 50000])

    out = tmp_path / "kyc.xlsx"
    wb.save(str(out))

    docs = await load_xlsx(out)
    assert docs[0]["metadata"]["header"] == [
        "aadhaar_last4", "pan_masked", "amount",
    ]
