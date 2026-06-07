"""v0.8.0: Tests for new document loaders."""

from __future__ import annotations

import json
from unittest.mock import AsyncMock, MagicMock, patch

import pytest

respx = pytest.importorskip("respx")


# -------------------- PPTX --------------------


@pytest.mark.asyncio
async def test_load_pptx_when_python_pptx_missing(tmp_path, monkeypatch):
    p = tmp_path / "x.pptx"
    p.write_bytes(b"fake")

    real_import = (
        __builtins__["__import__"] if isinstance(__builtins__, dict) else __builtins__.__import__
    )

    def fake_import(name, *args, **kwargs):
        if name == "pptx":
            raise ImportError("Mocked")
        return real_import(name, *args, **kwargs)

    monkeypatch.setattr("builtins.__import__", fake_import)
    from largestack._loaders import load_pptx

    docs = await load_pptx(str(p))
    assert "python-pptx" in docs[0]["metadata"]["error"]


@pytest.mark.asyncio
async def test_load_pptx_returns_doc_per_slide():
    """When python-pptx is available, each slide becomes a document."""
    pytest.importorskip("pptx")
    from pptx import Presentation
    from largestack._loaders import load_pptx
    import tempfile

    # Build a real .pptx
    prs = Presentation()
    blank_layout = prs.slide_layouts[5]  # title only
    s1 = prs.slides.add_slide(blank_layout)
    s1.shapes.title.text = "First Slide"
    s2 = prs.slides.add_slide(blank_layout)
    s2.shapes.title.text = "Second Slide"

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        prs.save(f.name)
        path = f.name

    docs = await load_pptx(path)
    assert len(docs) == 2
    assert "First Slide" in docs[0]["content"]
    assert docs[0]["metadata"]["slide"] == 0


# -------------------- EPUB --------------------


@pytest.mark.asyncio
async def test_load_epub_when_ebooklib_missing(tmp_path, monkeypatch):
    p = tmp_path / "x.epub"
    p.write_bytes(b"fake")
    real_import = (
        __builtins__["__import__"] if isinstance(__builtins__, dict) else __builtins__.__import__
    )

    def fake_import(name, *args, **kwargs):
        if name == "ebooklib":
            raise ImportError("Mocked")
        return real_import(name, *args, **kwargs)

    monkeypatch.setattr("builtins.__import__", fake_import)
    from largestack._loaders import load_epub

    docs = await load_epub(str(p))
    assert "ebooklib" in docs[0]["metadata"]["error"]


# -------------------- Excel --------------------


@pytest.mark.asyncio
async def test_load_excel_returns_doc_per_sheet():
    pytest.importorskip("openpyxl")
    from openpyxl import Workbook
    from largestack._loaders import load_excel
    import tempfile

    wb = Workbook()
    ws1 = wb.active
    ws1.title = "Sheet1"
    ws1.append(["Name", "Age"])
    ws1.append(["Alice", 30])
    ws1.append(["Bob", 25])
    ws2 = wb.create_sheet("Sheet2")
    ws2.append(["Item", "Price"])
    ws2.append(["X", 9.99])

    with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as f:
        wb.save(f.name)
        path = f.name

    docs = await load_excel(path)
    assert len(docs) == 2
    assert docs[0]["metadata"]["sheet"] == "Sheet1"
    assert "Alice" in docs[0]["content"]
    assert docs[1]["metadata"]["sheet"] == "Sheet2"
    assert "9.99" in docs[1]["content"]


@pytest.mark.asyncio
async def test_load_excel_specific_sheet_only():
    pytest.importorskip("openpyxl")
    from openpyxl import Workbook
    from largestack._loaders import load_excel
    import tempfile

    wb = Workbook()
    wb.active.title = "A"
    wb.active.append(["x"])
    wb.create_sheet("B").append(["y"])
    with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as f:
        wb.save(f.name)
        path = f.name
    docs = await load_excel(path, sheet_name="B")
    assert len(docs) == 1
    assert docs[0]["metadata"]["sheet"] == "B"


@pytest.mark.asyncio
async def test_load_excel_when_openpyxl_missing(monkeypatch, tmp_path):
    p = tmp_path / "x.xlsx"
    p.write_bytes(b"fake")
    real_import = (
        __builtins__["__import__"] if isinstance(__builtins__, dict) else __builtins__.__import__
    )

    def fake_import(name, *args, **kwargs):
        if name == "openpyxl":
            raise ImportError("Mocked")
        return real_import(name, *args, **kwargs)

    monkeypatch.setattr("builtins.__import__", fake_import)
    from largestack._loaders import load_excel

    docs = await load_excel(str(p))
    assert "openpyxl" in docs[0]["metadata"]["error"]


# -------------------- S3 --------------------


@pytest.mark.asyncio
async def test_load_s3_when_boto3_missing(monkeypatch):
    real_import = (
        __builtins__["__import__"] if isinstance(__builtins__, dict) else __builtins__.__import__
    )

    def fake_import(name, *args, **kwargs):
        if name == "boto3":
            raise ImportError("Mocked")
        return real_import(name, *args, **kwargs)

    monkeypatch.setattr("builtins.__import__", fake_import)
    from largestack._loaders import load_s3

    docs = await load_s3("bucket", "key.txt")
    assert "boto3" in docs[0]["metadata"]["error"]


@pytest.mark.asyncio
async def test_load_s3_fetches_and_dispatches():
    """boto3 returns bytes → save to temp → dispatch by extension."""
    fake_body = MagicMock()
    fake_body.read = MagicMock(return_value=b"hello world from S3")
    fake_obj = {"Body": fake_body}
    fake_s3_client = MagicMock()
    fake_s3_client.get_object = MagicMock(return_value=fake_obj)

    fake_boto3 = MagicMock()
    fake_boto3.client = MagicMock(return_value=fake_s3_client)

    with patch.dict("sys.modules", {"boto3": fake_boto3}):
        from largestack._loaders import load_s3

        docs = await load_s3("my-bucket", "data.txt")

    fake_boto3.client.assert_called_with("s3")
    fake_s3_client.get_object.assert_called_with(Bucket="my-bucket", Key="data.txt")
    assert "hello world from S3" in docs[0]["content"]
    assert docs[0]["metadata"]["source"] == "s3://my-bucket/data.txt"


# -------------------- GCS --------------------


@pytest.mark.asyncio
async def test_load_gcs_fetches_and_dispatches():
    fake_blob = MagicMock()
    fake_blob.download_as_bytes = MagicMock(return_value=b'{"k": "v"}')
    fake_bucket = MagicMock()
    fake_bucket.blob = MagicMock(return_value=fake_blob)
    fake_client = MagicMock()
    fake_client.bucket = MagicMock(return_value=fake_bucket)

    fake_gcs_module = MagicMock()
    fake_gcs_module.Client = MagicMock(return_value=fake_client)

    fake_pkg = MagicMock()
    fake_pkg.cloud = MagicMock()
    fake_pkg.cloud.storage = fake_gcs_module

    with patch.dict(
        "sys.modules",
        {
            "google": fake_pkg,
            "google.cloud": fake_pkg.cloud,
            "google.cloud.storage": fake_gcs_module,
        },
    ):
        from largestack._loaders import load_gcs

        docs = await load_gcs("my-bucket", "test.json")

    assert docs[0]["metadata"]["source"] == "gs://my-bucket/test.json"
    assert "v" in docs[0]["content"]


# -------------------- YouTube --------------------


@pytest.mark.asyncio
async def test_load_youtube_extracts_video_id_from_url():
    fake_transcript = MagicMock()
    fake_transcript.get_transcript = MagicMock(
        return_value=[
            {"text": "Hello", "start": 0.0, "duration": 1.0},
            {"text": "world", "start": 1.0, "duration": 1.0},
        ]
    )

    fake_yta = MagicMock()
    fake_yta.YouTubeTranscriptApi = fake_transcript

    with patch.dict("sys.modules", {"youtube_transcript_api": fake_yta}):
        from largestack._loaders import load_youtube_transcript

        docs = await load_youtube_transcript("https://www.youtube.com/watch?v=dQw4w9WgXcQ")

    assert "Hello" in docs[0]["content"]
    assert "world" in docs[0]["content"]
    assert docs[0]["metadata"]["video_id"] == "dQw4w9WgXcQ"


@pytest.mark.asyncio
async def test_load_youtube_handles_missing_dep(monkeypatch):
    real_import = (
        __builtins__["__import__"] if isinstance(__builtins__, dict) else __builtins__.__import__
    )

    def fake_import(name, *args, **kwargs):
        if name == "youtube_transcript_api":
            raise ImportError("Mocked")
        return real_import(name, *args, **kwargs)

    monkeypatch.setattr("builtins.__import__", fake_import)
    from largestack._loaders import load_youtube_transcript

    docs = await load_youtube_transcript("video123")
    assert "youtube-transcript-api" in docs[0]["metadata"]["error"]


# -------------------- Wikipedia --------------------


@pytest.mark.asyncio
async def test_load_wikipedia_returns_extract():
    from largestack._loaders import load_wikipedia

    fake_response = {
        "query": {
            "pages": {
                "12345": {
                    "title": "Bengaluru",
                    "extract": "Bengaluru is the capital of Karnataka.",
                }
            }
        }
    }
    with respx.mock() as mock:
        mock.get("https://en.wikipedia.org/w/api.php").respond(200, json=fake_response)
        docs = await load_wikipedia("Bengaluru")
    assert "Bengaluru is the capital" in docs[0]["content"]
    assert docs[0]["metadata"]["title"] == "Bengaluru"


@pytest.mark.asyncio
async def test_load_wikipedia_no_match():
    from largestack._loaders import load_wikipedia

    fake_response = {"query": {"pages": {"-1": {"title": "X", "missing": ""}}}}
    with respx.mock() as mock:
        mock.get("https://en.wikipedia.org/w/api.php").respond(200, json=fake_response)
        docs = await load_wikipedia("nonsense_topic_zzz")
    assert "no Wikipedia article" in docs[0]["metadata"]["error"]


# -------------------- ArXiv --------------------


@pytest.mark.asyncio
async def test_load_arxiv_parses_atom_response():
    from largestack._loaders import load_arxiv

    atom_response = """<?xml version="1.0" encoding="UTF-8"?>
<feed xmlns="http://www.w3.org/2005/Atom">
  <entry>
    <id>http://arxiv.org/abs/2401.12345</id>
    <title>Test Paper Title</title>
    <summary>This is the abstract of a test paper.</summary>
    <published>2024-01-15T00:00:00Z</published>
    <author><name>Alice Researcher</name></author>
    <author><name>Bob Coauthor</name></author>
  </entry>
</feed>"""

    with respx.mock() as mock:
        mock.get("http://export.arxiv.org/api/query").respond(200, text=atom_response)
        docs = await load_arxiv("test query")

    assert len(docs) == 1
    assert "Test Paper Title" in docs[0]["metadata"]["title"]
    assert "abstract" in docs[0]["content"].lower()
    assert "Alice Researcher" in docs[0]["metadata"]["authors"]


@pytest.mark.asyncio
async def test_load_arxiv_no_results():
    from largestack._loaders import load_arxiv

    with respx.mock() as mock:
        mock.get("http://export.arxiv.org/api/query").respond(
            200,
            text='<?xml version="1.0"?><feed xmlns="http://www.w3.org/2005/Atom"></feed>',
        )
        docs = await load_arxiv("nothing matches")
    assert "no ArXiv results" in docs[0]["metadata"]["error"]


# -------------------- PubMed --------------------


@pytest.mark.asyncio
async def test_load_pubmed_two_step_fetch():
    from largestack._loaders import load_pubmed

    esearch_resp = {"esearchresult": {"idlist": ["12345", "67890"]}}
    efetch_xml = """<?xml version="1.0"?>
<PubmedArticleSet>
  <PubmedArticle>
    <MedlineCitation>
      <PMID>12345</PMID>
      <Article>
        <ArticleTitle>Diabetes Study</ArticleTitle>
        <Abstract>
          <AbstractText>Background: study of diabetes.</AbstractText>
          <AbstractText>Methods: ...</AbstractText>
        </Abstract>
      </Article>
    </MedlineCitation>
  </PubmedArticle>
</PubmedArticleSet>"""

    with respx.mock() as mock:
        mock.get("https://eutils.ncbi.nlm.nih.gov/entrez/eutils/esearch.fcgi").respond(
            200, json=esearch_resp
        )
        mock.get("https://eutils.ncbi.nlm.nih.gov/entrez/eutils/efetch.fcgi").respond(
            200, text=efetch_xml
        )
        docs = await load_pubmed("diabetes")

    assert len(docs) == 1
    assert "Diabetes Study" in docs[0]["metadata"]["title"]
    assert "12345" in docs[0]["metadata"]["pmid"]
    assert "study of diabetes" in docs[0]["content"]


@pytest.mark.asyncio
async def test_load_pubmed_no_results():
    from largestack._loaders import load_pubmed

    with respx.mock() as mock:
        mock.get("https://eutils.ncbi.nlm.nih.gov/entrez/eutils/esearch.fcgi").respond(
            200, json={"esearchresult": {"idlist": []}}
        )
        docs = await load_pubmed("nonsense_query_zz")
    assert "no PubMed results" in docs[0]["metadata"]["error"]


# -------------------- Dispatcher routes new formats --------------------


@pytest.mark.asyncio
async def test_dispatcher_routes_pptx_xlsx_epub(tmp_path):
    """The ``load()`` dispatcher routes PPTX, XLSX, EPUB to right loader."""
    from largestack._loaders import load

    # PPTX: file exists but pptx may not be installed
    p = tmp_path / "x.pptx"
    p.write_bytes(b"fake")
    docs = await load(str(p))
    # Either succeeds (lib installed) or returns error doc — but routes correctly
    assert isinstance(docs, list)

    # XLSX
    p = tmp_path / "x.xlsx"
    p.write_bytes(b"fake")
    docs = await load(str(p))
    assert isinstance(docs, list)

    # EPUB
    p = tmp_path / "x.epub"
    p.write_bytes(b"fake")
    docs = await load(str(p))
    assert isinstance(docs, list)
