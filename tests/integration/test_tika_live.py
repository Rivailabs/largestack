"""Optional live Apache Tika integration tests.

Run with:

    LARGESTACK_TIKA_LIVE=1 TIKA_SERVER_URL=http://127.0.0.1:9998 pytest tests/integration/test_tika_live.py -q
"""

from __future__ import annotations

import os

import pytest


pytestmark = pytest.mark.skipif(
    os.environ.get("LARGESTACK_TIKA_LIVE") != "1" or not os.environ.get("TIKA_SERVER_URL"),
    reason="set LARGESTACK_TIKA_LIVE=1 and TIKA_SERVER_URL to run live Tika tests",
)


@pytest.mark.asyncio
async def test_live_tika_text_and_html(tmp_path):
    from largestack._loaders.tika import load_with_tika

    text_path = tmp_path / "sample.txt"
    text_path.write_text("Largestack Tika text fixture", encoding="utf-8")

    html_path = tmp_path / "sample.html"
    html_path.write_text(
        "<html><body><h1>Largestack Tika HTML fixture</h1></body></html>",
        encoding="utf-8",
    )

    for path in (text_path, html_path):
        docs = await load_with_tika(path, fallback_on_error=False)
        assert docs
        assert "Largestack Tika" in "\n".join(doc["content"] for doc in docs)


@pytest.mark.asyncio
async def test_live_tika_generated_common_formats(tmp_path):
    from largestack._loaders.tika import load_with_tika

    paths = [_write_minimal_pdf(tmp_path / "sample.pdf")]

    try:
        import docx  # type: ignore
    except ImportError:
        docx = None
    if docx is not None:
        docx_path = tmp_path / "sample.docx"
        document = docx.Document()
        document.add_paragraph("Largestack Tika DOCX fixture")
        document.save(docx_path)
        paths.append(docx_path)

    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        Presentation = None
    if Presentation is not None:
        pptx_path = tmp_path / "sample.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = "Largestack Tika PPTX fixture"
        presentation.save(pptx_path)
        paths.append(pptx_path)

    try:
        from openpyxl import Workbook  # type: ignore
    except ImportError:
        Workbook = None
    if Workbook is not None:
        xlsx_path = tmp_path / "sample.xlsx"
        workbook = Workbook()
        sheet = workbook.active
        sheet["A1"] = "Largestack Tika XLSX fixture"
        workbook.save(xlsx_path)
        paths.append(xlsx_path)

    for path in paths:
        docs = await load_with_tika(path, fallback_on_error=False)
        assert docs
        text = "\n".join(doc["content"] for doc in docs)
        assert "Largestack Tika" in text


def _write_minimal_pdf(path):
    path.write_bytes(
        b"%PDF-1.4\n"
        b"1 0 obj << /Type /Catalog /Pages 2 0 R >> endobj\n"
        b"2 0 obj << /Type /Pages /Kids [3 0 R] /Count 1 >> endobj\n"
        b"3 0 obj << /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >> endobj\n"
        b"4 0 obj << /Length 71 >> stream\n"
        b"BT /F1 18 Tf 72 720 Td (Largestack Tika PDF fixture) Tj ET\n"
        b"endstream endobj\n"
        b"5 0 obj << /Type /Font /Subtype /Type1 /BaseFont /Helvetica >> endobj\n"
        b"xref\n0 6\n0000000000 65535 f \n"
        b"0000000009 00000 n \n0000000058 00000 n \n0000000115 00000 n \n"
        b"0000000241 00000 n \n0000000322 00000 n \n"
        b"trailer << /Size 6 /Root 1 0 R >>\nstartxref\n391\n%%EOF\n"
    )
    return path
